import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def extract_text_from_file(file_path: str) -> str:
    """
    Extract plain text from PDF, DOCX, or PPTX file.
    """
    extension = os.path.splitext(file_path)[1].lower()

    if extension == '.pdf':
        return _extract_pdf(file_path)

    elif extension == '.docx':
        return _extract_docx(file_path)

    elif extension == '.pptx':
        return _extract_pptx(file_path)

    else:
        raise ValueError("Unsupported file type")


def _extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)

    return "\n".join(text)


def _extract_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def _extract_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)
